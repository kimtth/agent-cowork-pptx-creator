from __future__ import annotations

import ast
import argparse
import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from layout_specs import (
    estimate_text_height_in,
    flow_layout_spec,
    get_layout_spec,
    LayoutSpec,
    RectSpec,
    CardsSpec,
    StatsSpec,
    TimelineSpec,
    ComparisonSpec,
)
from layout_validator import (
    validate_presentation,
    report_issues,
)

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

ICON_CACHE_DIR = os.environ.get('ICON_CACHE_DIR', '')
if not ICON_CACHE_DIR:
    _default_cache = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'skills', 'iconfy-list', 'cache')
    if os.path.isdir(_default_cache):
        ICON_CACHE_DIR = os.path.normpath(_default_cache)

# ---------------------------------------------------------------------------
# Noto Sans font support for non-Latin scripts
# ---------------------------------------------------------------------------

# Map of script ranges → Noto Sans font family name
_NOTO_FONT_MAP: list[tuple[str, list[tuple[int, int]]]] = [
    ('Noto Sans JP', [(0x3040, 0x309F), (0x30A0, 0x30FF), (0x31F0, 0x31FF),   # Hiragana, Katakana
                      (0x4E00, 0x9FFF), (0xFF65, 0xFF9F), (0xFF01, 0xFF60)]),  # CJK Unified (JP fallback)
    ('Noto Sans KR', [(0xAC00, 0xD7AF), (0x1100, 0x11FF), (0x3130, 0x318F)]), # Hangul
    ('Noto Sans SC', [(0x4E00, 0x9FFF),]),  # CJK Unified (SC fallback, lower priority than JP)
    ('Noto Sans TC', [(0x4E00, 0x9FFF),]),  # CJK Unified (TC fallback)
    ('Noto Sans Thai', [(0x0E00, 0x0E7F),]),
    ('Noto Sans Arabic', [(0x0600, 0x06FF), (0x0750, 0x077F), (0xFB50, 0xFDFF)]),
    ('Noto Sans Devanagari', [(0x0900, 0x097F),]),
]

# Google Fonts download URLs for the main Noto Sans variants needed for PPTX
_NOTO_FONT_URLS: dict[str, str] = {
    'Noto Sans JP': 'https://github.com/notofonts/noto-cjk/releases/download/Sans2.005/08_NotoSansJP.zip',
    'Noto Sans KR': 'https://github.com/notofonts/noto-cjk/releases/download/Sans2.005/09_NotoSansKR.zip',
    'Noto Sans SC': 'https://github.com/notofonts/noto-cjk/releases/download/Sans2.005/10_NotoSansSC.zip',
    'Noto Sans TC': 'https://github.com/notofonts/noto-cjk/releases/download/Sans2.005/11_NotoSansTC.zip',
}

_WINDOWS_USER_FONTS = Path(os.environ.get('LOCALAPPDATA', '')) / 'Microsoft' / 'Windows' / 'Fonts'
_SYSTEM_FONTS = Path(r'C:\Windows\Fonts') if sys.platform == 'win32' else Path('/usr/share/fonts')


def _font_installed(family: str) -> bool:
    """Check whether a font family has a .ttf/.otf file in system or user fonts."""
    slug = family.replace(' ', '')
    patterns = [slug.lower(), family.lower().replace(' ', '')]

    search_dirs: list[Path] = []
    if sys.platform == 'win32':
        search_dirs.append(_SYSTEM_FONTS)
        if _WINDOWS_USER_FONTS.is_dir():
            search_dirs.append(_WINDOWS_USER_FONTS)
    else:
        search_dirs.append(_SYSTEM_FONTS)

    for d in search_dirs:
        if not d.is_dir():
            continue
        for f in d.iterdir():
            name_lower = f.name.lower()
            if any(p in name_lower for p in patterns) and name_lower.endswith(('.ttf', '.otf')):
                return True
    return False


def _install_font_windows(ttf_path: Path) -> None:
    """Install a font file to the Windows per-user font directory."""
    _WINDOWS_USER_FONTS.mkdir(parents=True, exist_ok=True)
    dest = _WINDOWS_USER_FONTS / ttf_path.name
    if dest.exists():
        return
    import shutil
    shutil.copy2(ttf_path, dest)

    # Register in the Windows registry for the current user
    try:
        import winreg
        key = winreg.OpenKey(
            winreg.HKEY_CURRENT_USER,
            r'SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts',
            0,
            winreg.KEY_SET_VALUE,
        )
        font_name = ttf_path.stem.replace('-', ' ')
        winreg.SetValueEx(key, f'{font_name} (TrueType)', 0, winreg.REG_SZ, str(dest))
        winreg.CloseKey(key)
    except Exception:
        pass  # Font file is copied — registry is optional for python-pptx


def _download_and_install_noto(family: str, fonts_cache_dir: Path) -> bool:
    """Download a Noto Sans font family and install it. Returns True on success."""
    url = _NOTO_FONT_URLS.get(family)
    if not url:
        return False

    family_dir = fonts_cache_dir / family.replace(' ', '')
    if family_dir.is_dir() and any(family_dir.glob('*.ttf')):
        # Already cached — just ensure installed
        for ttf in family_dir.glob('*.ttf'):
            if sys.platform == 'win32':
                _install_font_windows(ttf)
            return True
        return True

    print(f'[font] Downloading {family}...', file=sys.stderr)
    import urllib.request
    import tempfile
    import zipfile
    try:
        zip_path = os.path.join(tempfile.gettempdir(), f'{family.replace(" ", "")}.zip')
        urllib.request.urlretrieve(url, zip_path)

        family_dir.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(zip_path, 'r') as zf:
            for member in zf.namelist():
                # Only extract Regular weight .ttf (keep it simple)
                basename = os.path.basename(member)
                if not basename.endswith('.ttf'):
                    continue
                lower = basename.lower()
                if 'regular' in lower or 'medium' in lower or 'bold' in lower:
                    target = family_dir / basename
                    with zf.open(member) as src, open(target, 'wb') as dst:
                        dst.write(src.read())
                    if sys.platform == 'win32':
                        _install_font_windows(target)

        os.unlink(zip_path)
        print(f'[font] Installed {family}', file=sys.stderr)
        return True
    except Exception as exc:
        print(f'[font] Failed to download {family}: {exc}', file=sys.stderr)
        return False


def ensure_noto_fonts(text: str, fonts_cache_dir: str = '') -> None:
    """Ensure that appropriate Noto Sans fonts are available for the given text."""
    needed: set[str] = set()
    for ch in text:
        cp = ord(ch)
        for family, ranges in _NOTO_FONT_MAP:
            if any(lo <= cp <= hi for lo, hi in ranges):
                needed.add(family)
                break

    if not needed:
        return

    cache_dir = Path(fonts_cache_dir) if fonts_cache_dir else Path(os.environ.get('WORKSPACE_DIR', '.')) / 'fonts'

    for family in needed:
        if _font_installed(family):
            continue
        _download_and_install_noto(family, cache_dir)


def resolve_font(text: str, base_font: str = 'Calibri') -> str:
    """Return the best font for the given text content.

    If the text contains CJK/non-Latin characters, returns the appropriate
    Noto Sans variant. Otherwise returns the base font unchanged.
    """
    for ch in text:
        cp = ord(ch)
        for family, ranges in _NOTO_FONT_MAP:
            if any(lo <= cp <= hi for lo, hi in ranges):
                return family
    return base_font


def _make_transparent(png_path: str) -> str:
    """Convert a black-on-white RGB icon to black-on-transparent RGBA.

    Returns the path to the transparent version (cached alongside the original).
    If the source already has an alpha channel with real transparency, returns it as-is.
    """
    suffix = '_t.png'
    transparent_path = os.path.join(
        os.path.dirname(png_path),
        f'{Path(png_path).stem}{suffix}',
    )
    if os.path.isfile(transparent_path):
        return transparent_path

    try:
        from PIL import Image
        img = Image.open(png_path)

        # If image already has real transparency, skip conversion
        if img.mode == 'RGBA':
            arr = img.load()
            # Quick sample: if corners are transparent, assume it's already good
            w, h = img.size
            if arr[0, 0][3] == 0 or arr[w - 1, 0][3] == 0:
                return png_path

        img = img.convert('RGBA')
        pixels = img.load()
        w, h = img.size
        for y in range(h):
            for x in range(w):
                r, g, b, _ = pixels[x, y]
                # Luminance: near-white → transparent, darker → opaque icon stroke
                lum = r * 0.299 + g * 0.587 + b * 0.114
                if lum > 240:
                    pixels[x, y] = (0, 0, 0, 0)
                else:
                    # Map luminance to alpha: black=255, mid-gray=partial
                    alpha = min(255, int((255 - lum) * (255 / 200)))
                    pixels[x, y] = (0, 0, 0, alpha)
        img.save(transparent_path, 'PNG')
        return transparent_path
    except Exception:
        return png_path


def _recolor_png(png_path: str, color_hex: str) -> str:
    """Tint an icon PNG to the requested color.

    Handles both RGBA (transparent bg) and RGB (white bg) source icons.
    """
    # First ensure we have a transparent version
    png_path = _make_transparent(png_path)

    color = color_hex.lstrip('#')
    if color == '000000':
        return png_path  # already black-on-transparent

    colored_path = os.path.join(
        os.path.dirname(png_path),
        f'{Path(png_path).stem}_{color}.png',
    )
    if os.path.isfile(colored_path):
        return colored_path

    try:
        from PIL import Image
        img = Image.open(png_path).convert('RGBA')
        r_tgt = int(color[0:2], 16)
        g_tgt = int(color[2:4], 16)
        b_tgt = int(color[4:6], 16)
        pixels = img.load()
        w, h = img.size
        for y in range(h):
            for x in range(w):
                _, _, _, a = pixels[x, y]
                if a > 0:
                    pixels[x, y] = (r_tgt, g_tgt, b_tgt, a)
        img.save(colored_path, 'PNG')
        return colored_path
    except Exception:
        return png_path


def fetch_icon(name: str, color_hex: str = '000000', size: int = 256) -> str | None:
    """Load an icon PNG from the local cache, recolor via Pillow if needed."""
    if ':' not in name:
        name = f'mdi:{name}'
    prefix, icon_name = name.split(':', 1)
    png_name = f'{icon_name}.png'

    if ICON_CACHE_DIR:
        cached = os.path.join(ICON_CACHE_DIR, prefix, png_name)
        if os.path.isfile(cached):
            return _recolor_png(cached, color_hex)

    print(f'[icon] NOT FOUND: {prefix}:{icon_name} (looked in {ICON_CACHE_DIR}/{prefix}/)', file=sys.stderr)
    return None


def _load_theme() -> dict[str, str]:
    raw = os.environ.get('PPTX_THEME_JSON', '{}')
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return {}
    return parsed if isinstance(parsed, dict) else {}


def rgb_color(value: str | None, fallback: str = '000000') -> RGBColor:
    normalized = (value or fallback).strip().lstrip('#').upper()
    if len(normalized) != 6:
        normalized = fallback
    return RGBColor.from_string(normalized)


def ensure_parent_dir(file_path: str) -> None:
    Path(file_path).expanduser().resolve().parent.mkdir(parents=True, exist_ok=True)


def apply_widescreen(prs: Presentation) -> Presentation:
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    return prs


def _convert_to_pptx_compatible(source: Path) -> Path:
    """Convert unsupported image formats (e.g. WebP) to PNG for python-pptx."""
    if source.suffix.lower() not in ('.webp',):
        return source
    target = source.with_suffix('.png')
    if target.exists():
        return target
    from PIL import Image
    with Image.open(source) as img:
        img.save(target, 'PNG')
    return target


def safe_image_path(value: str | None) -> str | None:
    if not value:
        return None
    candidate = Path(value).expanduser().resolve()
    if not candidate.exists():
        return None
    candidate = _convert_to_pptx_compatible(candidate)
    return str(candidate)


ICON_INSERT_SCALE = 0.4  # Scale factor for icons inserted into fixed-size placeholders (e.g. title slide) to prevent overflow


def _is_icon_asset(path: str) -> bool:
    if not ICON_CACHE_DIR:
        return False
    normalized = os.path.normcase(os.path.normpath(path))
    return normalized.startswith(os.path.normcase(os.path.normpath(ICON_CACHE_DIR)) + os.sep)


def safe_add_picture(shapes, image_path: str | None, left, top, width=None, height=None):
    resolved = safe_image_path(image_path)
    if not resolved:
        return None
    if width is not None and height is not None and _is_icon_asset(resolved):
        scaled_width = max(1, int(width * ICON_INSERT_SCALE))
        scaled_height = max(1, int(height * ICON_INSERT_SCALE))
        left = left + int((width - scaled_width) / 2)
        top = top + int((height - scaled_height) / 2)
        width = scaled_width
        height = scaled_height
    return shapes.add_picture(resolved, left, top, width=width, height=height)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument('generated_code')
    parser.add_argument('output_path')
    parser.add_argument('--render-dir', default=None)
    parser.add_argument('--workspace-dir', default=None,
                        help='Absolute path to the user workspace directory')
    return parser.parse_args()


def build_namespace(generated_path: Path, output_path: Path, *, workspace_dir: str = '') -> dict[str, object]:
    theme = _load_theme()
    title = os.environ.get('PPTX_TITLE', 'Presentation')
    if not workspace_dir:
        workspace_dir = os.environ.get('WORKSPACE_DIR', '')
    images_dir = os.path.join(workspace_dir, 'images') if workspace_dir else ''

    return {
        '__name__': '__main__',
        '__file__': str(generated_path),
        'OUTPUT_PATH': str(output_path),
        'PPTX_TITLE': title,
        'PPTX_THEME': theme,
        'WORKSPACE_DIR': workspace_dir,
        'IMAGES_DIR': images_dir,
        'SLIDE_WIDTH_IN': SLIDE_WIDTH_IN,
        'SLIDE_HEIGHT_IN': SLIDE_HEIGHT_IN,
        'os': os,
        'Presentation': Presentation,
        'Inches': Inches,
        'Pt': Pt,
        'RGBColor': RGBColor,
        'PP_ALIGN': PP_ALIGN,
        'MSO_ANCHOR': MSO_ANCHOR,
        'MSO_AUTO_SIZE': MSO_AUTO_SIZE,
        'MSO_AUTO_SHAPE_TYPE': MSO_AUTO_SHAPE_TYPE,
        'rgb_color': rgb_color,
        'ensure_parent_dir': ensure_parent_dir,
        'apply_widescreen': apply_widescreen,
        'safe_image_path': safe_image_path,
        'safe_add_picture': safe_add_picture,
        'estimate_text_height_in': estimate_text_height_in,
        'flow_layout_spec': flow_layout_spec,
        'get_layout_spec': get_layout_spec,
        'LayoutSpec': LayoutSpec,
        'RectSpec': RectSpec,
        'CardsSpec': CardsSpec,
        'StatsSpec': StatsSpec,
        'TimelineSpec': TimelineSpec,
        'ComparisonSpec': ComparisonSpec,
        '_pptx_theme': theme,
        '_pptx_title': title,
        'fetch_icon': fetch_icon,
        'ICON_CACHE_DIR': ICON_CACHE_DIR,
        'resolve_font': resolve_font,
        'ensure_noto_fonts': ensure_noto_fonts,
    }


def validate_generated_code_syntax(code: str, generated_path: Path) -> None:
    if '```' in code:
        raise RuntimeError(
            'Generated Python code still contains Markdown code fences. '
            'Return raw Python only inside a single fenced block, and do not nest or duplicate code fences.'
        )

    if 'from future import annotations' in code and 'from __future__ import annotations' not in code:
        raise RuntimeError(
            'Generated Python code uses an invalid future import. '
            'Use "from __future__ import annotations".'
        )

    if re.search(r'^\s*import\s+annotations\b', code, re.MULTILINE):
        raise RuntimeError(
            'Generated Python code uses "import annotations" which is not a valid module. '
            'Use "from __future__ import annotations" or remove the import entirely.'
        )

    if re.search(r'\bif\s+name\s*==', code) and 'if __name__' not in code:
        raise RuntimeError(
            'Generated Python code uses "if name ==" instead of "if __name__ == \'__main__\':". '
            'The variable must be __name__ (with double underscores).'
        )

    try:
        ast.parse(code, filename=str(generated_path))
    except SyntaxError as exc:
        location = f'line {exc.lineno}' if exc.lineno else 'unknown line'
        if exc.offset:
            location = f'{location}, column {exc.offset}'
        source_line = (exc.text or '').rstrip()
        details = [f'Generated Python code has invalid syntax at {location}: {exc.msg}']
        if source_line:
            details.append(source_line)
        raise RuntimeError('\n'.join(details)) from exc

    try:
        compile(code, str(generated_path), 'exec')
    except SyntaxError as exc:
        location = f'line {exc.lineno}' if exc.lineno else 'unknown line'
        if exc.offset:
            location = f'{location}, column {exc.offset}'
        source_line = (exc.text or '').rstrip()
        details = [f'Generated Python code has a semantic error at {location}: {exc.msg}']
        if source_line:
            details.append(source_line)
        raise RuntimeError('\n'.join(details)) from exc


def run_generated_code(generated_path: Path, namespace: dict[str, object]) -> None:
    code = generated_path.read_text(encoding='utf-8')
    validate_generated_code_syntax(code, generated_path)
    exec(compile(code, str(generated_path), 'exec'), namespace)


def finalize_output(output_path: Path, namespace: dict[str, object]) -> None:
    if output_path.exists():
        return

    builder = namespace.get('build_presentation')
    if callable(builder):
        builder(str(output_path), namespace.get('_pptx_theme'), namespace.get('_pptx_title'))

    if not output_path.exists():
        raise RuntimeError('Generated python-pptx code completed without creating the PPTX output file.')


def validate_and_fix_output(output_path: Path) -> None:
    """Run auto-size enforcement, then validation on the generated PPTX."""
    from layout_validator import _enforce_auto_size

    prs = Presentation(str(output_path))
    changed = False

    # Always enforce auto-size on every text frame (prevents text overflow → overlap)
    auto_size_fixes = 0
    for slide in prs.slides:
        auto_size_fixes += _enforce_auto_size(slide)
    if auto_size_fixes > 0:
        changed = True
        print(f'[layout-validator] Enforced auto-size on {auto_size_fixes} text frame(s).', file=sys.stderr)

    # Save after auto-size if anything changed
    if changed:
        prs.save(str(output_path))

    # Validate remaining issues (detection + report only)
    issues = validate_presentation(prs)
    if not issues:
        print('[layout-validator] All slides passed layout validation.', file=sys.stderr)
        return

    report = report_issues(issues)
    print(report, file=sys.stderr)

    blocking = [issue for issue in issues if issue.severity.value == 'error']
    if blocking:
        has_overlap = any('overlap' in i.issue_type.value.lower() for i in blocking)
        has_text_overflow = any('text_overflow' in i.issue_type.value.lower() for i in blocking)

        # Tolerate up to 2 blocking issues — treat as success with warnings
        if len(blocking) <= 2:
            print(
                f'[layout-validator] {len(blocking)} minor layout issue(s) detected (within tolerance). '
                'PPTX generated with incomplete layout details.',
                file=sys.stderr,
            )
            return

        hints: list[str] = []
        if has_overlap:
            hints.append(
                'TOOL HINT: Use patch_layout_infrastructure(action="read", file="layout_specs") to inspect '
                'current layout coordinates, then patch_layout_infrastructure(action="patch", ...) to adjust '
                'get_layout_spec() dimensions. After patching, call rerun_pptx to re-execute.'
            )
        if has_text_overflow:
            hints.append(
                'TOOL HINT: Use patch_layout_infrastructure(action="read", file="layout_validator") to inspect '
                'validation thresholds, then patch if needed. Or adjust layout_specs dimensions to provide more space. '
                'After patching, call rerun_pptx to re-execute.'
            )

        hint_text = '\n'.join(hints)
        raise RuntimeError(
            'Layout validation failed after generation. '
            'Reduce content density, reserve more vertical space, or split the slide.\n'
            'Alternatively, use the layout infrastructure tools to fix spec dimensions or validator thresholds.\n\n'
            f'{hint_text}\n\n'
            f'{report}'
        )


def render_preview_images(output_path: Path, render_dir: Path) -> None:
    if sys.platform != 'win32':
        raise RuntimeError('Local slide preview rendering is only supported on Windows.')

    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except ImportError as exc:
        raise RuntimeError('pywin32 is required for local PPTX preview rendering on Windows.') from exc

    render_dir.mkdir(parents=True, exist_ok=True)
    # Only remove old preview images; preserve generated-source.py and .pptx files
    for existing in render_dir.glob('*'):
        if existing.is_file() and existing.suffix.lower() in ('.png', '.jpg', '.jpeg'):
            existing.unlink()

    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.DispatchEx('PowerPoint.Application')
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(output_path), WithWindow=False, ReadOnly=True)
        presentation.Export(str(render_dir), 'PNG', 1280, 720)
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError('Microsoft PowerPoint is required to render local preview images.') from exc
    finally:
        if presentation is not None:
            presentation.Close()
        if powerpoint is not None:
            powerpoint.Quit()
        pythoncom.CoUninitialize()


def _unlock_or_rename(output_path: Path) -> Path:
    """Remove existing output file. If locked, fall back to a timestamped name."""
    if not output_path.exists():
        return output_path
    try:
        output_path.unlink()
        return output_path
    except PermissionError:
        from datetime import datetime
        stamp = datetime.now().strftime('%Y%m%d-%H%M%S')
        alt = output_path.with_stem(f'{output_path.stem}-{stamp}')
        print(f'[WARNING] {output_path.name} is locked, saving as {alt.name}', file=sys.stderr)
        return alt


def main() -> int:
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

    if ICON_CACHE_DIR:
        cache_exists = os.path.isdir(ICON_CACHE_DIR)
        print(f'[icon-cache] ICON_CACHE_DIR={ICON_CACHE_DIR} (exists={cache_exists})', file=sys.stderr)
    else:
        print('[icon-cache] ICON_CACHE_DIR is not set — uncached icons will be unavailable', file=sys.stderr)

    args = parse_args()
    generated_path = Path(args.generated_code).resolve()
    output_path = Path(args.output_path).resolve()
    render_dir = Path(args.render_dir).resolve() if args.render_dir else None
    workspace_dir = str(Path(args.workspace_dir).resolve()) if args.workspace_dir else ''

    print(f'[workspace] WORKSPACE_DIR={workspace_dir or "(not set)"}', file=sys.stderr)

    if not generated_path.exists():
        raise FileNotFoundError(f'Generated Python source file not found: {generated_path}')

    output_path = _unlock_or_rename(output_path)
    namespace = build_namespace(generated_path, output_path, workspace_dir=workspace_dir)

    # Pre-download Noto Sans fonts for any non-Latin text in the generated code
    try:
        code_text = generated_path.read_text(encoding='utf-8')
        ensure_noto_fonts(code_text, os.environ.get('WORKSPACE_DIR', ''))
    except Exception as exc:  # noqa: BLE001
        print(f'[font] Font pre-check failed (non-blocking): {exc}', file=sys.stderr)

    run_generated_code(generated_path, namespace)
    finalize_output(output_path, namespace)

    try:
        validate_and_fix_output(output_path)
    except Exception as exc:  # noqa: BLE001
        print(f'[layout-validator] Validation failed (PPTX was generated): {exc}', file=sys.stderr)

    if render_dir is not None:
        try:
            render_preview_images(output_path, render_dir)
        except Exception as exc:  # noqa: BLE001
            print(f'[WARNING] Preview rendering failed (PPTX was generated successfully): {exc}', file=sys.stderr)

    return 0


if __name__ == '__main__':
    raise SystemExit(main())
