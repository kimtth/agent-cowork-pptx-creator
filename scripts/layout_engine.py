"""Post-generation layout engine for python-pptx slides.

After the LLM generates a PPTX file, this engine re-positions shapes to
eliminate overlaps using a greedy top-to-bottom vertical packer.

This is the REPAIR layer — it processes shapes in display order (top-to-bottom)
within column lanes and pushes shapes downward when overlaps are detected.
Unlike pairwise nudging, it never creates new overlaps because shapes are
packed sequentially.

Usage:
    from layout_engine import relayout_presentation

    prs = Presentation('deck.pptx')
    relayout_presentation(prs)
    prs.save('deck.pptx')
"""

from __future__ import annotations

from dataclasses import dataclass
from pptx.util import Inches

SLIDE_WIDTH_EMU = Inches(13.333)
SLIDE_HEIGHT_EMU = Inches(7.5)
MIN_GAP_EMU = Inches(0.12)

# Shapes smaller than this threshold are considered decorations
# and won't be repositioned independently
DECORATION_MAX_AREA_RATIO = 0.008  # 0.8% of slide area
SLIDE_AREA_EMU2 = SLIDE_WIDTH_EMU * SLIDE_HEIGHT_EMU

# Shapes covering more than this fraction of the slide are backgrounds
BACKGROUND_COVER_RATIO = 0.90


@dataclass
class _ShapeInfo:
    """Internal representation of a shape for layout calculations."""
    shape: object  # python-pptx shape
    left: int
    top: int
    width: int
    height: int
    name: str
    is_decoration: bool
    parent_index: int | None  # index of parent shape, if decoration

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    @property
    def area(self) -> int:
        return self.width * self.height

    @property
    def cx(self) -> float:
        return self.left + self.width / 2.0

    @property
    def cy(self) -> float:
        return self.top + self.height / 2.0


def _extract_shapes(slide) -> list[_ShapeInfo]:
    """Extract positionable shapes from a slide."""
    infos: list[_ShapeInfo] = []
    for shape in slide.shapes:
        try:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
        except (AttributeError, TypeError):
            continue
        if left is None or top is None or width is None or height is None:
            continue
        if width <= 0 or height <= 0:
            continue

        name = getattr(shape, 'name', '') or ''

        # Skip notes/footer shapes — these are protected from repositioning
        name_lower = name.lower()
        if name_lower.startswith('notes') or name_lower.startswith('footer'):
            continue

        area = width * height

        # Skip background fills
        is_bg = (width >= SLIDE_WIDTH_EMU * BACKGROUND_COVER_RATIO and
                 height >= SLIDE_HEIGHT_EMU * BACKGROUND_COVER_RATIO)
        if is_bg:
            continue

        is_dec = area < SLIDE_AREA_EMU2 * DECORATION_MAX_AREA_RATIO
        infos.append(_ShapeInfo(
            shape=shape,
            left=left,
            top=top,
            width=width,
            height=height,
            name=name,
            is_decoration=is_dec,
            parent_index=None,
        ))
    return infos


def _find_parent(dec: _ShapeInfo, shapes: list[_ShapeInfo]) -> int | None:
    """Find the closest non-decoration shape that contains or nearly contains this decoration."""
    best_index = None
    best_distance = float('inf')

    for i, s in enumerate(shapes):
        if s.is_decoration:
            continue
        # Check if decoration center is within the shape bounds (with tolerance)
        tol = Inches(0.3)
        if (s.left - tol <= dec.cx <= s.right + tol and
                s.top - tol <= dec.cy <= s.bottom + tol):
            dist = abs(dec.cx - s.cx) + abs(dec.cy - s.cy)
            if dist < best_distance:
                best_distance = dist
                best_index = i
    return best_index


def _shapes_overlap_significantly(a: _ShapeInfo, b: _ShapeInfo) -> bool:
    """Check if two shapes overlap by more than a trivial amount."""
    ox = max(0, min(a.right, b.right) - max(a.left, b.left))
    oy = max(0, min(a.bottom, b.bottom) - max(a.top, b.top))
    if ox <= 0 or oy <= 0:
        return False
    overlap_area = ox * oy
    smaller_area = min(a.area, b.area)
    if smaller_area <= 0:
        return False
    return overlap_area / smaller_area > 0.08


def _assign_lane(shape: _ShapeInfo) -> str:
    """Assign shape to a column lane based on horizontal position."""
    mid_x = SLIDE_WIDTH_EMU * 0.65
    # Wide shapes (>65% of slide width) are in the "full" lane
    if shape.width > SLIDE_WIDTH_EMU * 0.65:
        return 'full'
    if shape.left + shape.width / 2 < mid_x:
        return 'left'
    return 'right'


def relayout_slide(slide) -> int:
    """Re-layout shapes on a single slide to eliminate overlaps.

    Returns the number of shapes repositioned.
    """
    shapes = _extract_shapes(slide)
    if len(shapes) < 2:
        return 0

    # Step 1: Link decorations to parent shapes
    for i, s in enumerate(shapes):
        if s.is_decoration:
            s.parent_index = _find_parent(s, shapes)

    # Step 2: Separate independent shapes from decorations
    independent = [s for s in shapes if not s.is_decoration or s.parent_index is None]
    decorations = [(i, s) for i, s in enumerate(shapes) if s.is_decoration and s.parent_index is not None]

    # Step 3: Group independent shapes by lane
    lanes: dict[str, list[_ShapeInfo]] = {'full': [], 'left': [], 'right': []}
    for s in independent:
        lane = _assign_lane(s)
        lanes[lane].append(s)

    # Step 4: Sort each lane by top position (display order)
    for lane_shapes in lanes.values():
        lane_shapes.sort(key=lambda s: (s.top, s.left))

    repositioned = 0

    # Step 5: Greedy vertical packing within each lane
    for lane_name, lane_shapes in lanes.items():
        if len(lane_shapes) < 2:
            continue

        # Track the occupied vertical ranges for this lane
        placed: list[_ShapeInfo] = [lane_shapes[0]]

        for s in lane_shapes[1:]:
            # Check overlap with all previously placed shapes in this lane
            need_push = False
            push_below_y = s.top

            for placed_s in placed:
                if _shapes_overlap_significantly(s, placed_s):
                    need_push = True
                    candidate_y = placed_s.bottom + int(MIN_GAP_EMU)
                    push_below_y = max(push_below_y, candidate_y)

            # Also check cross-lane overlaps with 'full' width shapes
            if lane_name != 'full':
                for full_s in lanes['full']:
                    if _shapes_overlap_significantly(s, full_s):
                        need_push = True
                        candidate_y = full_s.bottom + int(MIN_GAP_EMU)
                        push_below_y = max(push_below_y, candidate_y)

            if need_push and push_below_y != s.top:
                # Ensure we don't push beyond slide bottom
                new_top = min(push_below_y, int(SLIDE_HEIGHT_EMU) - s.height)
                if new_top > s.top:
                    # Verify the clamped position doesn't still collide with placed shapes.
                    # This can happen when multiple shapes overflow and all get clamped
                    # to the same slide-bottom Y — moving them would just stack them.
                    test = _ShapeInfo(
                        shape=None, left=s.left, top=new_top,
                        width=s.width, height=s.height,
                        name=s.name, is_decoration=False, parent_index=None,
                    )
                    still_blocked = any(_shapes_overlap_significantly(test, p) for p in placed)
                    if not still_blocked:
                        delta = new_top - s.top
                        s.shape.top = new_top
                        s.top = new_top  # update our tracking
                        repositioned += 1

                        # Move associated decorations by the same delta
                        for _, dec in decorations:
                            if dec.parent_index is not None:
                                parent = shapes[dec.parent_index]
                                if parent is s:
                                    dec.shape.top = dec.top + delta
                                    dec.top += delta

            placed.append(s)

    # Step 6: Cross-lane overlap check (left vs right)
    for left_s in lanes['left']:
        for right_s in lanes['right']:
            if _shapes_overlap_significantly(left_s, right_s):
                # Push the right-lane shape right if horizontally overlapping
                if right_s.left < left_s.right:
                    new_left = left_s.right + int(MIN_GAP_EMU)
                    if new_left + right_s.width <= int(SLIDE_WIDTH_EMU):
                        right_s.shape.left = new_left
                        right_s.left = new_left
                        repositioned += 1

    return repositioned


def relayout_presentation(prs) -> int:
    """Re-layout all slides in a presentation.

    Returns total number of shapes repositioned.
    """
    total = 0
    for slide in prs.slides:
        total += relayout_slide(slide)
    return total
