import os
import json
import math
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def build_presentation(output_path, theme, title):
    prs = apply_widescreen(Presentation())

    # ── Load layout input ──
    layout_input_path = os.path.join(WORKSPACE_DIR, 'previews', 'layout-input.json')
    with open(layout_input_path, 'r', encoding='utf-8') as f:
        layout_input = json.load(f)

    # ── Theme color map ──
    c = {
        'bg': theme.get('BG', 'FCFCFD'),
        'text': theme.get('TEXT', '181F25'),
        'dark': theme.get('DARK', '181F25'),
        'dark2': theme.get('DARK2', '2B333B'),
        'light': theme.get('LIGHT', 'FCFCFD'),
        'light2': theme.get('LIGHT2', 'DEFCF0'),
        'accent1': theme.get('ACCENT1', 'A8FAD8'),
        'accent2': theme.get('ACCENT2', '75FFD1'),
        'accent3': theme.get('ACCENT3', 'A5F3FD'),
        'accent4': theme.get('ACCENT4', '5CE4FF'),
        'accent5': theme.get('ACCENT5', 'D8EBFD'),
        'accent6': theme.get('ACCENT6', 'B1D6FC'),
        'link': theme.get('LINK', 'FCFCFD'),
        'primary': theme.get('PRIMARY', 'A8FAD8'),
        'white': theme.get('WHITE', 'FCFCFD'),
        'border': theme.get('BORDER', 'E1E1E1'),
    }

    # ── Slide data ──
    slides_data = [
        {
            'title': 'Microsoft Copilot Studio',
            'key_message': 'AIエージェントで業務を変革する時代が来た',
            'bullets': ['ローコードで構築するインテリジェントAIエージェント',
                        'Microsoft 365 と Power Platform のエコシステム統合',
                        '導入ガイド 2026'],
            'icon': 'lucide:bot',
            'image': os.path.join(IMAGES_DIR, '01-microsoft-copilot-studio-d80c5481fa.jpg'),
        },
        {
            'title': 'エグゼクティブサマリー：3つの導入メリット',
            'key_message': 'Copilot Studioはローコード開発・エンタープライズ統合・堅牢なガバナンスを同時に実現する',
            'bullets': ['ローコードで迅速構築 — 開発者不要でAIエージェントを数時間で展開',
                        'エンタープライズ統合 — M365・Dynamics 365・外部APIとシームレス接続',
                        'セキュリティ＆ガバナンス — DLP・監査ログ・MIP感度ラベルで企業基準を維持'],
            'icon': 'lucide:sparkles',
            'image': None,
        },
        {
            'title': '現状：AI活用の課題が業務変革を阻んでいる',
            'key_message': '多くの組織がAI導入で開発リソース不足・統合の複雑さ・セキュリティ懸念に直面している',
            'bullets': ['AI開発には専門人材が必要 — データサイエンティスト・開発者の確保が困難',
                        '既存システムとの統合が複雑 — サイロ化したデータとプロセス',
                        'セキュリティとコンプライアンスの確保 — 生成AIのガバナンスが未整備',
                        'スケーラビリティの壁 — PoC から本番展開への移行が難航'],
            'icon': 'lucide:alert-triangle',
            'image': None,
        },
        {
            'title': 'Copilot Studioとは：ローコードでAIエージェントを構築するプラットフォーム',
            'key_message': 'エージェントとエージェントフローにより、会話型AI・タスク自動化・ワークフロー統合を一つのプラットフォームで実現',
            'bullets': ['エージェント — NLU＋生成AIで自律的に対話・判断・行動するAIコンパニオン',
                        'エージェントフロー — 自然言語またはビジュアルエディターで反復タスクを自動化',
                        'ナレッジ統合 — SharePoint・Graph・外部コネクタからリアルタイムに知識を取得',
                        'マルチチャネル展開 — Teams・Web・モバイル・WhatsAppなど複数チャネルに公開'],
            'icon': 'lucide:layers',
            'image': None,
        },
        {
            'title': 'ツール選択：Agent Builder と Copilot Studio を使い分ける',
            'key_message': '軽量Q&Aは M365 Agent Builder、高度なワークフロー・外部統合は Copilot Studio を選択すべき',
            'bullets': ['M365 Agent Builder — 個人・小規模チーム向け、自然言語オーサリング、SharePoint Q&A中心',
                        'Copilot Studio — 部門・組織・外部顧客向け、マルチステップワークフロー、Premium コネクタ',
                        '移行パス — Agent Builder で作成 → Copilot Studio にコピーして高度な機能を解放',
                        '判断基準 — 対象ユーザー規模・機能要件・ガバナンスニーズで選択'],
            'icon': 'lucide:git-branch',
            'image': None,
        },
        {
            'title': '進化し続ける：2025〜2026年の主要リリース',
            'key_message': 'GPT-5対応・コンピュータ使用エージェント・MCP統合など、半年ごとに大型機能が投入されている',
            'bullets': ['2025年9月 — Computer-Using Agents（CUA）プレビュー、Client SDK公開',
                        '2025年10月 — GPT-4.1標準化、MCPサーバー統合、エージェント評価テストセット',
                        '2025年11月 — GPT-5 GA展開、マルチエージェント連携、SharePointメタデータフィルター',
                        '2026年1月 — VS Code拡張GA、Cloud PCプーリング、評価機能強化',
                        '2026年2月 — Claude Sonnet 4.5対応、プロンプトビルダー強化',
                        '2026年3月 — Work IQツール統合（プレビュー）'],
            'icon': 'lucide:calendar',
            'image': None,
        },
        {
            'title': 'ライセンス体系：3つの導入パスから選べる',
            'key_message': '無料試用版→M365付属プラン→スタンドアロンサブスクリプションの段階的導入が可能',
            'bullets': ['無料試用版 — サインアップ即日開始、30日延長可、試用後90日間エージェント稼働継続',
                        'M365サブスクリプション付属 — Teams限定チャネル、標準コネクタ、基本機能',
                        'スタンドアロンサブスクリプション — 全チャネル展開、Premiumコネクタ、生成オーケストレーション、ALM',
                        '従量課金制 — Copilotクレジットベースの柔軟な課金オプション'],
            'icon': 'lucide:credit-card',
            'image': None,
        },
        {
            'title': 'エンタープライズガバナンス：セキュリティは組み込み済み',
            'key_message': 'DLP・監査ログ・CMK・MIPラベルなど、企業レベルのセキュリティ制御がプラットフォームに標準装備されている',
            'bullets': ['データ損失防止（DLP） — Power Platform データポリシーでコネクタ・チャネル・生成AIを制御',
                        '監査とモニタリング — Microsoft Purview＋Sentinel で作成者・ユーザー行動を完全可視化',
                        '暗号化とデータ所在地 — CMK対応、地理的データレジデンシー準拠',
                        '感度ラベル — MIPラベルがSharePointナレッジソースからエージェント応答まで一貫適用'],
            'icon': 'lucide:shield-check',
            'image': os.path.join(IMAGES_DIR, '08-image-fa181b767d.png'),
        },
        {
            'title': '推奨アクション：90日間の導入ロードマップ',
            'key_message': '試用版開始→PoC構築→本番展開の3フェーズで、リスクを最小化しながら迅速に導入できる',
            'bullets': ['Phase 1（1〜30日）— 無料試用版でサインアップ、IT部門Q&AエージェントでPoCを構築',
                        'Phase 2（31〜60日）— Agent Builder vs Copilot Studio の使い分けを確定、DLPポリシー設計',
                        'Phase 3（61〜90日）— スタンドアロンライセンス取得、本番チャネル展開、監査ログ有効化'],
            'icon': 'lucide:rocket',
            'image': None,
        },
        {
            'title': 'Next Steps',
            'key_message': '今日から始められる3つの具体的アクションがある',
            'bullets': ['① copilotstudio.microsoft.com で無料試用版にサインアップ',
                        '② 自社のユースケースを特定し、最初のエージェントPoCを構築',
                        '③ IT管理者と連携し、DLPポリシーと監査ログの設計を開始'],
            'icon': 'lucide:arrow-right',
            'image': None,
        },
    ]

    # ── Helper: add Swiss International left vertical accent bar ──
    def add_left_bar(slide, color_hex, x=0.0, y=0.0, w=0.12, h=7.5):
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb_color(color_hex)
        bar.line.fill.background()
        bar.name = 'swiss_left_bar'

    # ── Helper: add horizontal divider rule ──
    def add_hrule(slide, spec_accent, color_hex):
        if spec_accent is None:
            return
        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(spec_accent.x), Inches(spec_accent.y),
            Inches(spec_accent.w), Inches(spec_accent.h)
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = rgb_color(color_hex)
        rule.line.fill.background()
        rule.name = 'swiss_hrule'

    # ── Helper: set slide background ──
    def set_bg(slide, color_hex):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb_color(color_hex)

    # ── Helper: add text box ──
    def add_textbox(slide, rect, text, font_size, color_hex, bold=False,
                    align=PP_ALIGN.LEFT, auto_size_mode=MSO_AUTO_SIZE.NONE,
                    font_name=None, name=None, anchor=MSO_ANCHOR.TOP,
                    line_spacing=None):
        box = slide.shapes.add_textbox(
            Inches(rect.x), Inches(rect.y),
            Inches(rect.w), Inches(rect.h)
        )
        if name:
            box.name = name
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = auto_size_mode
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        resolved = resolve_font(text, font_name or 'Calibri')
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb_color(color_hex)
        run.font.name = resolved
        if line_spacing:
            p.line_spacing = line_spacing
        return box

    # ── Helper: add panel (rounded rect with text) ──
    def add_panel(slide, x, y, w, h, fill_hex, texts, font_sizes, colors_list,
                  bolds, name=None, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE):
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb_color(fill_hex)
        panel.line.fill.background()
        if name:
            panel.name = name
        # Subtle corner radius
        if hasattr(panel, 'adjustments') and len(panel.adjustments) > 0:
            panel.adjustments[0] = 0.04
        tf = panel.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        for idx, txt in enumerate(texts):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = txt
            p.alignment = PP_ALIGN.LEFT
            if line_spacing:
                p.line_spacing = line_spacing
            resolved = resolve_font(txt, 'Calibri')
            run = p.runs[0]
            run.font.size = Pt(font_sizes[idx] if idx < len(font_sizes) else font_sizes[-1])
            run.font.bold = bolds[idx] if idx < len(bolds) else False
            run.font.color.rgb = rgb_color(colors_list[idx] if idx < len(colors_list) else colors_list[-1])
            run.font.name = resolved
        return panel

    # ── Helper: add icon to slide ──
    def place_icon(slide, spec, icon_name, color_hex):
        icon_path = fetch_icon(icon_name, color_hex=color_hex)
        if icon_path and spec.icon_rect:
            safe_add_picture(slide.shapes, icon_path,
                Inches(spec.icon_rect.x), Inches(spec.icon_rect.y),
                width=Inches(spec.icon_rect.w), height=Inches(spec.icon_rect.h))

    # ── Color accents per slide for variety ──
    accent_cycle = [c['accent1'], c['accent4'], c['accent3'], c['accent2'],
                    c['accent6'], c['accent5'], c['accent1'], c['accent4'],
                    c['accent3'], c['accent2']]
    card_fill_cycle = [c['light2'], c['accent5'], c['accent3'], c['accent6'],
                       c['accent5'], c['light2']]

    # ================================================================
    # SLIDE 0: Title
    # ================================================================
    si = 0
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent1'])

    # Hero image in hero_rect
    if sd['image'] and spec.hero_rect:
        safe_add_picture(slide.shapes, sd['image'],
            Inches(spec.hero_rect.x), Inches(spec.hero_rect.y),
            width=Inches(spec.hero_rect.w), height=Inches(spec.hero_rect.h))

    # Accent bar
    add_hrule(slide, spec.accent_rect, c['accent1'])

    # Title
    add_textbox(slide, spec.title_rect, sd['title'], 36, c['dark'],
                bold=True, font_name='Calibri', name='title')

    # Key message
    add_textbox(slide, spec.key_message_rect, sd['key_message'], 20, c['dark2'],
                font_name='Calibri', name='key_message', line_spacing=1.4)

    # Chips (subtitle bullets)
    if spec.chips_rect and sd['bullets']:
        chip_count = len(sd['bullets'])
        chip_gap = 0.15
        total_gap = chip_gap * (chip_count - 1)
        chip_w = (spec.chips_rect.w - total_gap) / chip_count
        for idx, bullet in enumerate(sd['bullets']):
            cx = spec.chips_rect.x + idx * (chip_w + chip_gap)
            chip_fill = card_fill_cycle[idx % len(card_fill_cycle)]
            fg = ensure_contrast(c['dark2'], chip_fill)
            add_panel(slide, cx, spec.chips_rect.y, chip_w, spec.chips_rect.h,
                      chip_fill, [bullet], [12], [fg], [False],
                      name=f'chip_{idx}', anchor=MSO_ANCHOR.MIDDLE)

    # Footer text
    if spec.footer_rect:
        add_textbox(slide, spec.footer_rect,
                    'Microsoft Copilot Studio 導入ガイド 2026',
                    11, c['dark2'], font_name='Calibri', name='footer_info',
                    align=PP_ALIGN.LEFT)

    # Icon
    place_icon(slide, spec, sd['icon'], c['accent4'])

    # Notes
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 1: Cards — Executive Summary
    # ================================================================
    si = 1
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent2'])
    add_hrule(slide, spec.accent_rect, c['accent2'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 22, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 14, c['dark2'],
                font_name='Calibri', name='key_message')

    card_accents_s1 = [c['accent1'], c['accent4'], c['accent6']]
    if spec.cards:
        for idx, bullet in enumerate(sd['bullets']):
            cr = spec.cards.card_rect(idx)
            parts = bullet.split(' — ', 1)
            card_title = parts[0]
            card_body = parts[1] if len(parts) > 1 else ''
            accent_bar_color = card_accents_s1[idx % len(card_accents_s1)]
            # Card background
            add_panel(slide, cr.x, cr.y, cr.w, cr.h,
                      c['light2'],
                      [card_title, card_body] if card_body else [card_title],
                      [15, 13] if card_body else [15],
                      [ensure_contrast(c['dark'], c['light2']),
                       ensure_contrast(c['dark2'], c['light2'])],
                      [True, False],
                      name=f'card_{idx}')
            # Left accent stripe on card
            stripe = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cr.x), Inches(cr.y), Inches(0.08), Inches(cr.h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rgb_color(accent_bar_color)
            stripe.line.fill.background()

    place_icon(slide, spec, sd['icon'], c['accent2'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 2: Bullets — Current Challenges
    # ================================================================
    si = 2
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent3'])
    add_hrule(slide, spec.accent_rect, c['accent3'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 22, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 14, c['dark2'],
                font_name='Calibri', name='key_message')

    bullet_accents = [c['accent4'], c['accent3'], c['accent6'], c['accent5']]
    if spec.content_rect:
        bullet_count = len(sd['bullets'])
        gap = 0.16
        row_h = (spec.content_rect.h - gap * (bullet_count - 1)) / bullet_count
        row_h = min(row_h, 0.85)
        for idx, bullet in enumerate(sd['bullets']):
            by = spec.content_rect.y + idx * (row_h + gap)
            parts = bullet.split(' — ', 1)
            b_title = parts[0]
            b_body = parts[1] if len(parts) > 1 else ''
            # Bullet dot
            dot_size = 0.14
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(spec.content_rect.x + 0.1),
                Inches(by + row_h / 2 - dot_size / 2),
                Inches(dot_size), Inches(dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb_color(bullet_accents[idx % len(bullet_accents)])
            dot.line.fill.background()
            # Text
            text_x = spec.content_rect.x + 0.4
            text_w = spec.content_rect.w - 0.5
            add_panel(slide, text_x, by, text_w, row_h,
                      c['bg'],
                      [b_title, b_body] if b_body else [b_title],
                      [16, 14] if b_body else [16],
                      [c['dark'], c['dark2']],
                      [True, False],
                      name=f'bullet_{idx}', line_spacing=1.3)

    place_icon(slide, spec, sd['icon'], c['accent3'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 3: Cards — What is Copilot Studio
    # ================================================================
    si = 3
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent4'])
    add_hrule(slide, spec.accent_rect, c['accent4'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 20, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 13, c['dark2'],
                font_name='Calibri', name='key_message')

    card_fills_s3 = [c['light2'], c['accent5'], c['light2'], c['accent5']]
    card_accent_s3 = [c['accent1'], c['accent4'], c['accent3'], c['accent6']]
    if spec.cards:
        for idx, bullet in enumerate(sd['bullets']):
            cr = spec.cards.card_rect(idx)
            parts = bullet.split(' — ', 1)
            card_title = parts[0]
            card_body = parts[1] if len(parts) > 1 else ''
            fill = card_fills_s3[idx % len(card_fills_s3)]
            fg_title = ensure_contrast(c['dark'], fill)
            fg_body = ensure_contrast(c['dark2'], fill)
            add_panel(slide, cr.x, cr.y, cr.w, cr.h, fill,
                      [card_title, card_body] if card_body else [card_title],
                      [15, 13] if card_body else [15],
                      [fg_title, fg_body],
                      [True, False],
                      name=f'card_{idx}')
            stripe = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cr.x), Inches(cr.y), Inches(0.08), Inches(cr.h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rgb_color(card_accent_s3[idx % len(card_accent_s3)])
            stripe.line.fill.background()

    place_icon(slide, spec, sd['icon'], c['accent4'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 4: Comparison — Agent Builder vs Copilot Studio
    # ================================================================
    si = 4
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent6'])
    add_hrule(slide, spec.accent_rect, c['accent6'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 20, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 13, c['dark2'],
                font_name='Calibri', name='key_message')

    if spec.comparison:
        left = spec.comparison.left
        right = spec.comparison.right

        # Left column header
        lh_h = 0.5
        add_panel(slide, left.x, left.y, left.w, lh_h,
                  c['accent5'],
                  ['M365 Agent Builder'],
                  [16], [ensure_contrast(c['dark'], c['accent5'])],
                  [True], name='comp_left_header', anchor=MSO_ANCHOR.MIDDLE)

        # Right column header
        add_panel(slide, right.x, right.y, right.w, lh_h,
                  c['accent1'],
                  ['Copilot Studio'],
                  [16], [ensure_contrast(c['dark'], c['accent1'])],
                  [True], name='comp_right_header', anchor=MSO_ANCHOR.MIDDLE)

        # Left body
        left_bullets = sd['bullets'][:2]
        left_body_y = left.y + lh_h + 0.12
        left_body_h = left.h - lh_h - 0.12
        for idx, b in enumerate(left_bullets):
            parts = b.split(' — ', 1)
            body_text = parts[1] if len(parts) > 1 else parts[0]
            item_h = (left_body_h - 0.12) / 2
            iy = left_body_y + idx * (item_h + 0.12)
            add_panel(slide, left.x, iy, left.w, item_h,
                      c['light2'],
                      [body_text], [14],
                      [ensure_contrast(c['dark2'], c['light2'])],
                      [False], name=f'comp_left_{idx}')

        # Right body
        right_bullets = sd['bullets'][2:]
        right_body_y = right.y + lh_h + 0.12
        right_body_h = right.h - lh_h - 0.12
        for idx, b in enumerate(right_bullets):
            parts = b.split(' — ', 1)
            body_text = parts[1] if len(parts) > 1 else parts[0]
            item_h = (right_body_h - 0.12) / 2
            iy = right_body_y + idx * (item_h + 0.12)
            add_panel(slide, right.x, iy, right.w, item_h,
                      c['light2'],
                      [body_text], [14],
                      [ensure_contrast(c['dark2'], c['light2'])],
                      [False], name=f'comp_right_{idx}')

    place_icon(slide, spec, sd['icon'], c['accent6'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 5: Timeline — Release Roadmap
    # ================================================================
    si = 5
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent5'])
    add_hrule(slide, spec.accent_rect, c['accent5'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 20, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 13, c['dark2'],
                font_name='Calibri', name='key_message')

    if spec.timeline:
        tl = spec.timeline
        # Vertical timeline line
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(tl.line_x), Inches(tl.line_y),
            Inches(0.04), Inches(tl.line_h)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = rgb_color(c['accent1'])
        line.line.fill.background()

        tl_accents = [c['accent1'], c['accent4'], c['accent3'],
                      c['accent2'], c['accent6'], c['accent5']]

        for idx, bullet in enumerate(sd['bullets']):
            node = tl.node_rect(idx)
            # Timeline dot
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(tl.dot_x), Inches(node.y + 0.08),
                Inches(tl.dot_size), Inches(tl.dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb_color(tl_accents[idx % len(tl_accents)])
            dot.line.fill.background()

            # Split date and description
            parts = bullet.split(' — ', 1)
            date_text = parts[0]
            desc_text = parts[1] if len(parts) > 1 else ''

            # Date label
            date_box = slide.shapes.add_textbox(
                Inches(node.x), Inches(node.y),
                Inches(1.6), Inches(node.h)
            )
            dtf = date_box.text_frame
            dtf.word_wrap = True
            dtf.auto_size = MSO_AUTO_SIZE.NONE
            dp = dtf.paragraphs[0]
            dp.text = date_text
            dp.alignment = PP_ALIGN.LEFT
            dr = dp.runs[0]
            dr.font.size = Pt(13)
            dr.font.bold = True
            dr.font.color.rgb = rgb_color(tl_accents[idx % len(tl_accents)])
            dr.font.name = resolve_font(date_text, 'Calibri')

            # Description
            if desc_text:
                desc_box = slide.shapes.add_textbox(
                    Inches(node.x + 1.7), Inches(node.y),
                    Inches(node.w - 1.7), Inches(node.h)
                )
                desc_tf = desc_box.text_frame
                desc_tf.word_wrap = True
                desc_tf.auto_size = MSO_AUTO_SIZE.NONE
                desc_p = desc_tf.paragraphs[0]
                desc_p.text = desc_text
                desc_p.alignment = PP_ALIGN.LEFT
                desc_r = desc_p.runs[0]
                desc_r.font.size = Pt(13)
                desc_r.font.color.rgb = rgb_color(c['dark2'])
                desc_r.font.name = resolve_font(desc_text, 'Calibri')

    place_icon(slide, spec, sd['icon'], c['accent5'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 6: Stats — Licensing
    # ================================================================
    si = 6
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent1'])
    add_hrule(slide, spec.accent_rect, c['accent1'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 22, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 14, c['dark2'],
                font_name='Calibri', name='key_message')

    stat_labels = ['無料試用版', 'M365付属', 'スタンドアロン']
    stat_numbers = ['0円', 'M365内', 'Premium']
    stat_descs = [
        'サインアップ即日開始\n30日延長可\n試用後90日間稼働継続',
        'Teams限定チャネル\n標準コネクタ\n基本機能',
        '全チャネル展開\nPremiumコネクタ\n生成オーケストレーション / ALM',
    ]
    stat_accents = [c['accent1'], c['accent4'], c['accent6']]

    if spec.stats:
        for idx in range(3):
            br = spec.stats.box_rect(idx)
            # Stat box background
            fill = c['light2']
            panel = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(br.x), Inches(br.y), Inches(br.w), Inches(br.h)
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = rgb_color(fill)
            panel.line.fill.background()
            if hasattr(panel, 'adjustments') and len(panel.adjustments) > 0:
                panel.adjustments[0] = 0.04
            panel.name = f'stat_box_{idx}'

            # Top accent stripe
            stripe = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(br.x), Inches(br.y), Inches(br.w), Inches(0.06)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rgb_color(stat_accents[idx])
            stripe.line.fill.background()

            # Number
            num_box = slide.shapes.add_textbox(
                Inches(br.x + 0.15), Inches(br.y + 0.2),
                Inches(br.w - 0.3), Inches(0.6)
            )
            ntf = num_box.text_frame
            ntf.word_wrap = True
            ntf.auto_size = MSO_AUTO_SIZE.NONE
            np_ = ntf.paragraphs[0]
            np_.text = stat_numbers[idx]
            np_.alignment = PP_ALIGN.CENTER
            nr = np_.runs[0]
            nr.font.size = Pt(32)
            nr.font.bold = True
            nr.font.color.rgb = rgb_color(stat_accents[idx])
            nr.font.name = resolve_font(stat_numbers[idx], 'Calibri')

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(br.x + 0.15), Inches(br.y + 0.8),
                Inches(br.w - 0.3), Inches(0.4)
            )
            ltf = label_box.text_frame
            ltf.word_wrap = True
            ltf.auto_size = MSO_AUTO_SIZE.NONE
            lp = ltf.paragraphs[0]
            lp.text = stat_labels[idx]
            lp.alignment = PP_ALIGN.CENTER
            lr = lp.runs[0]
            lr.font.size = Pt(16)
            lr.font.bold = True
            lr.font.color.rgb = rgb_color(ensure_contrast(c['dark'], fill))
            lr.font.name = resolve_font(stat_labels[idx], 'Calibri')

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(br.x + 0.15), Inches(br.y + 1.3),
                Inches(br.w - 0.3), Inches(br.h - 1.5)
            )
            dtf = desc_box.text_frame
            dtf.word_wrap = True
            dtf.auto_size = MSO_AUTO_SIZE.NONE
            for line_idx, line in enumerate(stat_descs[idx].split('\n')):
                if line_idx == 0:
                    dp = dtf.paragraphs[0]
                else:
                    dp = dtf.add_paragraph()
                dp.text = line
                dp.alignment = PP_ALIGN.CENTER
                dp.line_spacing = 1.4
                dr = dp.runs[0]
                dr.font.size = Pt(12)
                dr.font.color.rgb = rgb_color(ensure_contrast(c['dark2'], fill))
                dr.font.name = resolve_font(line, 'Calibri')

    # Footer with 4th bullet
    if spec.footer_rect:
        add_panel(slide, spec.footer_rect.x, spec.footer_rect.y,
                  spec.footer_rect.w, spec.footer_rect.h,
                  c['accent5'],
                  ['従量課金制 — Copilotクレジットベースの柔軟な課金オプション'],
                  [13], [ensure_contrast(c['dark2'], c['accent5'])],
                  [False], name='footer_note', anchor=MSO_ANCHOR.MIDDLE)

    place_icon(slide, spec, sd['icon'], c['accent1'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 7: Cards — Governance
    # ================================================================
    si = 7
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent4'])
    add_hrule(slide, spec.accent_rect, c['accent4'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 20, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 13, c['dark2'],
                font_name='Calibri', name='key_message')

    card_fills_s7 = [c['light2'], c['accent5'], c['light2'], c['accent5']]
    card_accents_s7 = [c['accent4'], c['accent1'], c['accent6'], c['accent3']]
    if spec.cards:
        for idx, bullet in enumerate(sd['bullets']):
            cr = spec.cards.card_rect(idx)
            parts = bullet.split(' — ', 1)
            card_title = parts[0]
            card_body = parts[1] if len(parts) > 1 else ''
            fill = card_fills_s7[idx % len(card_fills_s7)]
            fg_t = ensure_contrast(c['dark'], fill)
            fg_b = ensure_contrast(c['dark2'], fill)
            add_panel(slide, cr.x, cr.y, cr.w, cr.h, fill,
                      [card_title, card_body] if card_body else [card_title],
                      [15, 13] if card_body else [15],
                      [fg_t, fg_b],
                      [True, False],
                      name=f'card_{idx}')
            stripe = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cr.x), Inches(cr.y), Inches(0.08), Inches(cr.h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rgb_color(card_accents_s7[idx % len(card_accents_s7)])
            stripe.line.fill.background()

    # Image in icon area (slide 8 has an attached image)
    if sd['image'] and spec.icon_rect:
        safe_add_picture(slide.shapes, sd['image'],
            Inches(spec.icon_rect.x), Inches(spec.icon_rect.y),
            width=Inches(spec.icon_rect.w), height=Inches(spec.icon_rect.h))
    else:
        place_icon(slide, spec, sd['icon'], c['accent4'])

    # Also place the icon if image was placed
    if sd['image']:
        icon_path = fetch_icon(sd['icon'], color_hex=c['accent4'])
        if icon_path:
            safe_add_picture(slide.shapes, icon_path,
                Inches(spec.icon_rect.x + spec.icon_rect.w - 0.6),
                Inches(spec.icon_rect.y + spec.icon_rect.h - 0.6),
                width=Inches(0.5), height=Inches(0.5))

    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 8: Bullets — 90-day Roadmap
    # ================================================================
    si = 8
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, c['bg'])
    add_left_bar(slide, c['accent3'])
    add_hrule(slide, spec.accent_rect, c['accent3'])

    add_textbox(slide, spec.title_rect, sd['key_message'], 22, c['dark'],
                bold=True, font_name='Calibri', name='title')
    add_textbox(slide, spec.key_message_rect, sd['title'], 14, c['dark2'],
                font_name='Calibri', name='key_message')

    phase_accents = [c['accent1'], c['accent4'], c['accent6']]
    phase_labels = ['Phase 1', 'Phase 2', 'Phase 3']
    if spec.content_rect:
        bullet_count = len(sd['bullets'])
        gap = 0.2
        row_h = (spec.content_rect.h - gap * (bullet_count - 1)) / bullet_count
        row_h = min(row_h, 1.1)
        for idx, bullet in enumerate(sd['bullets']):
            by = spec.content_rect.y + idx * (row_h + gap)
            parts = bullet.split('— ', 1)
            phase_title = parts[0].strip()
            phase_body = parts[1].strip() if len(parts) > 1 else ''

            # Phase number box
            phase_box_w = 1.2
            add_panel(slide, spec.content_rect.x, by, phase_box_w, row_h,
                      phase_accents[idx],
                      [phase_labels[idx]],
                      [16], [ensure_contrast(c['dark'], phase_accents[idx])],
                      [True], name=f'phase_num_{idx}', anchor=MSO_ANCHOR.MIDDLE)

            # Phase content
            text_x = spec.content_rect.x + phase_box_w + 0.15
            text_w = spec.content_rect.w - phase_box_w - 0.15
            add_panel(slide, text_x, by, text_w, row_h,
                      c['light2'],
                      [phase_body] if phase_body else [phase_title],
                      [14],
                      [ensure_contrast(c['dark2'], c['light2'])],
                      [False], name=f'phase_body_{idx}')

    place_icon(slide, spec, sd['icon'], c['accent3'])
    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ================================================================
    # SLIDE 9: Summary — Next Steps
    # ================================================================
    si = 9
    spec = PRECOMPUTED_LAYOUT_SPECS[si]
    sd = slides_data[si]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Dark background for closing impact
    set_bg(slide, c['dark'])
    add_left_bar(slide, c['accent1'], h=7.5)

    # Title
    add_textbox(slide, spec.title_rect, sd['title'], 30, c['accent1'],
                bold=True, font_name='Calibri', name='title')
    # Key message
    add_textbox(slide, spec.key_message_rect, sd['key_message'], 16,
                ensure_contrast(c['light'], c['dark']),
                font_name='Calibri', name='key_message')

    # Summary box
    if spec.summary_box:
        # Accent line above summary
        accent_line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(spec.summary_box.x), Inches(spec.summary_box.y - 0.08),
            Inches(2.0), Inches(0.04)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = rgb_color(c['accent1'])
        accent_line.line.fill.background()

    # Action items
    action_accents = [c['accent1'], c['accent4'], c['accent6']]
    if spec.content_rect:
        bullet_count = len(sd['bullets'])
        gap = 0.2
        row_h = (spec.content_rect.h - gap * (bullet_count - 1)) / bullet_count
        row_h = min(row_h, 0.85)
        for idx, bullet in enumerate(sd['bullets']):
            by = spec.content_rect.y + idx * (row_h + gap)

            # Number circle
            circle_size = 0.4
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(spec.content_rect.x + 0.1),
                Inches(by + row_h / 2 - circle_size / 2),
                Inches(circle_size), Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = rgb_color(action_accents[idx])
            circle.line.fill.background()
            ctf = circle.text_frame
            ctf.word_wrap = False
            cp = ctf.paragraphs[0]
            cp.text = str(idx + 1)
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.runs[0]
            cr.font.size = Pt(16)
            cr.font.bold = True
            cr.font.color.rgb = rgb_color(ensure_contrast(c['dark'], action_accents[idx]))
            cr.font.name = 'Calibri'

            # Action text (strip the leading number)
            action_text = bullet.lstrip('①②③ ')
            text_x = spec.content_rect.x + 0.7
            text_w = spec.content_rect.w - 0.8
            action_box = slide.shapes.add_textbox(
                Inches(text_x), Inches(by),
                Inches(text_w), Inches(row_h)
            )
            atf = action_box.text_frame
            atf.word_wrap = True
            atf.auto_size = MSO_AUTO_SIZE.NONE
            ap = atf.paragraphs[0]
            ap.text = action_text
            ap.alignment = PP_ALIGN.LEFT
            resolved = resolve_font(action_text, 'Calibri')
            ar = ap.runs[0]
            ar.font.size = Pt(16)
            ar.font.color.rgb = rgb_color(ensure_contrast(c['light'], c['dark']))
            ar.font.name = resolved

    place_icon(slide, spec, sd['icon'], c['accent1'])

    # Notes citation
    if spec.notes_rect:
        notes_box = slide.shapes.add_textbox(
            Inches(spec.notes_rect.x), Inches(spec.notes_rect.y),
            Inches(spec.notes_rect.w), Inches(spec.notes_rect.h)
        )
        notes_box.name = 'notes_body'
        ntf = notes_box.text_frame
        ntf.word_wrap = True
        ntf.auto_size = MSO_AUTO_SIZE.NONE
        np_ = ntf.paragraphs[0]
        np_.text = 'Source: Microsoft Learn — Copilot Studio Documentation 2026'
        np_.alignment = PP_ALIGN.LEFT
        nr = np_.runs[0]
        nr.font.size = Pt(8)
        nr.font.color.rgb = rgb_color(c['accent5'])
        nr.font.name = 'Calibri'

    slide.notes_slide.notes_text_frame.text = layout_input[si].get('notes', '')

    # ── Save ──
    prs.save(output_path)