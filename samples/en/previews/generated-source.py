import os
import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(output_path, theme, title):
    prs = apply_widescreen(Presentation())

    # --- Load layout input ---
    layout_input_path = os.path.join(WORKSPACE_DIR, 'previews', 'layout-input.json')
    with open(layout_input_path, 'r', encoding='utf-8') as f:
        layout_input = json.load(f)

    # --- Theme color mapping ---
    c = {
        'bg': theme.get('BG', 'FCFCFD'),
        'text': theme.get('TEXT', '181F25'),
        'dark': theme.get('DARK', '181F25'),
        'dark2': theme.get('DARK2', '2B333B'),
        'light': theme.get('LIGHT', 'FCFCFD'),
        'light2': theme.get('LIGHT2', 'DEFCF0'),
        'accent1': theme.get('ACCENT1', 'A8FAD8'),
        'accent2': theme.get('ACCENT2', '75FFD1'),
        'accent3': theme.get('ACCENT3', 'A5F3FD'),
        'accent4': theme.get('ACCENT4', '5CE4FF'),
        'accent5': theme.get('ACCENT5', 'D8EBFD'),
        'accent6': theme.get('ACCENT6', 'B1D6FC'),
        'primary': theme.get('PRIMARY', 'A8FAD8'),
        'white': theme.get('WHITE', 'FCFCFD'),
        'border': theme.get('BORDER', 'E1E1E1'),
    }

    # Swiss International accent colors mapped to theme
    ACCENT_BAR_COLOR = c['primary']
    DIVIDER_COLOR = c['border']
    CARD_ACCENTS = [c['accent1'], c['accent3'], c['accent4'], c['accent5'], c['accent6'], c['accent2']]
    SOFT_BG = c['light2']

    FONT_TITLE = 'Calibri'
    FONT_BODY = 'Calibri'

    def rc(hex_str):
        return rgb_color(hex_str)

    def add_vertical_accent_bar(slide, spec, color_hex=None):
        """Swiss International signature: vertical bar on left edge."""
        bar_color = color_hex or ACCENT_BAR_COLOR
        bar_w = 0.08
        bar_x = 0.22
        bar_y = 0.5
        bar_h = 6.5
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(bar_x), Inches(bar_y),
            Inches(bar_w), Inches(bar_h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rc(bar_color)
        bar.line.fill.background()
        bar.rotation = 0.0

    def add_horizontal_divider(slide, y, x=0.5, w=12.333, color_hex=None):
        """Swiss International horizontal rule."""
        div_color = color_hex or DIVIDER_COLOR
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = rc(div_color)
        line.line.fill.background()

    def add_slide_number(slide, slide_num):
        """Small slide number in bottom-right corner."""
        txBox = slide.shapes.add_textbox(
            Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3)
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(slide_num)
        run.font.size = Pt(9)
        run.font.color.rgb = rc(c['dark2'])
        run.font.name = FONT_BODY

    def add_title_textbox(slide, spec, text, font_size=28, bold=True):
        """Add title text using spec.title_rect."""
        r = spec.title_rect
        txBox = slide.shapes.add_textbox(
            Inches(r.x), Inches(r.y), Inches(r.w), Inches(r.h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rc(c['text'])
        run.font.name = resolve_font(text, FONT_TITLE)
        return txBox

    def add_key_message(slide, spec, text, font_size=16):
        """Add key message using spec.key_message_rect."""
        r = spec.key_message_rect
        if r is None:
            return None
        txBox = slide.shapes.add_textbox(
            Inches(r.x), Inches(r.y), Inches(r.w), Inches(r.h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = False
        run.font.color.rgb = rc(c['dark2'])
        run.font.name = resolve_font(text, FONT_BODY)
        return txBox

    def add_accent_divider(slide, spec, color_hex=None):
        """Add accent divider bar from spec.accent_rect."""
        r = spec.accent_rect
        if r is None:
            return
        add_horizontal_divider(slide, r.y, r.x, r.w, color_hex or ACCENT_BAR_COLOR)

    def add_icon_to_slide(slide, spec, icon_name, color_hex):
        """Add icon using spec.icon_rect."""
        icon_path = fetch_icon(icon_name, color_hex=color_hex)
        if icon_path and spec.icon_rect:
            ir = spec.icon_rect
            safe_add_picture(slide.shapes, icon_path,
                Inches(ir.x), Inches(ir.y),
                width=Inches(ir.w), height=Inches(ir.h))

    def add_notes(slide, notes_text):
        """Add speaker notes."""
        slide.notes_slide.notes_text_frame.text = notes_text

    def set_slide_bg(slide, color_hex=None):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = rc(color_hex or c['bg'])

    def write_panel_text(shape, title_text, body_text, title_pt=15, body_pt=14,
                         title_color=None, body_color=None, accent_color=None):
        """Write title + body into a panel shape with auto-fit."""
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.08)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(title_pt)
        run.font.bold = True
        run.font.color.rgb = rc(title_color or c['text'])
        run.font.name = resolve_font(title_text, FONT_TITLE)

        if body_text:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(4)
            p2.alignment = PP_ALIGN.LEFT
            run2 = p2.add_run()
            run2.text = body_text
            run2.font.size = Pt(body_pt)
            run2.font.bold = False
            run2.font.color.rgb = rc(body_color or c['dark2'])
            run2.font.name = resolve_font(body_text, FONT_BODY)

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide_idx = 0
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Hero image in right zone
    hero = spec.hero_rect
    if hero:
        img_path = os.path.join(IMAGES_DIR,
            '01-microsoft-copilot-studio-transform-operations-with-low-code--d80c5481fa.jpg')
        safe_add_picture(slide.shapes, img_path,
            Inches(hero.x), Inches(hero.y),
            width=Inches(hero.w), height=Inches(hero.h))

    # Left accent bar
    add_vertical_accent_bar(slide, spec)

    # Title
    add_title_textbox(slide, spec, sd['title_text'], font_size=30, bold=True)

    # Key message
    add_key_message(slide, spec, sd['key_message_text'], font_size=17)

    # Accent divider
    add_accent_divider(slide, spec)

    # Chips (3 bullets as horizontal labels)
    chips = spec.chips_rect
    if chips and sd.get('bullets'):
        chip_count = len(sd['bullets'])
        gap = 0.15
        total_gap = gap * (chip_count - 1)
        chip_w = (chips.w - total_gap) / chip_count
        for ci, bullet in enumerate(sd['bullets']):
            cx = chips.x + ci * (chip_w + gap)
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(chips.y),
                Inches(chip_w), Inches(chips.h)
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = rc(SOFT_BG)
            chip.line.color.rgb = rc(c['accent1'])
            chip.line.width = Pt(1)
            tf = chip.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(11)
            run.font.color.rgb = rc(c['text'])
            run.font.name = resolve_font(bullet, FONT_BODY)

    # Icon
    add_icon_to_slide(slide, spec, 'lucide:bot', c['accent4'])

    add_slide_number(slide, 1)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 2: Cards — Conclusion
    # =========================================================================
    slide_idx = 1
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=28)
    add_key_message(slide, spec, sd['key_message_text'], font_size=16)
    add_accent_divider(slide, spec)

    # 3 cards in a row
    cards = spec.cards
    if cards:
        card_colors = [c['accent1'], c['accent3'], c['accent5']]
        for ci, bullet in enumerate(sd['bullets'][:3]):
            parts = bullet.split(' — ', 1)
            card_title = parts[0] if len(parts) > 1 else bullet
            card_body = parts[1] if len(parts) > 1 else ''
            cr = cards.card_rect(ci)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cr.x), Inches(cr.y),
                Inches(cr.w), Inches(cr.h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = rc(c['white'])
            card.line.color.rgb = rc(card_colors[ci % len(card_colors)])
            card.line.width = Pt(2)

            # Left accent stripe on card
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(cr.x), Inches(cr.y),
                Inches(0.06), Inches(cr.h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rc(card_colors[ci % len(card_colors)])
            stripe.line.fill.background()

            write_panel_text(card, card_title, card_body, title_pt=15, body_pt=13)

    add_icon_to_slide(slide, spec, 'lucide:sparkles', c['accent2'])
    add_slide_number(slide, 2)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 3: Bullets — Reality Check
    # =========================================================================
    slide_idx = 2
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=28)
    add_key_message(slide, spec, sd['key_message_text'], font_size=16)
    add_accent_divider(slide, spec)

    # 4 bullet rows as styled blocks
    content = spec.content_rect
    if content and sd.get('bullets'):
        bullet_count = len(sd['bullets'])
        gap = 0.12
        row_h = (content.h - gap * (bullet_count - 1)) / bullet_count
        row_h = min(row_h, 0.72)
        bullet_icons = ['lucide:users', 'lucide:unplug', 'lucide:shield-alert', 'lucide:trending-up']

        for bi, bullet in enumerate(sd['bullets'][:4]):
            parts = bullet.split(' — ', 1)
            b_title = parts[0]
            b_body = parts[1] if len(parts) > 1 else ''
            by = content.y + bi * (row_h + gap)

            # Bullet row background
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(content.x + 0.15), Inches(by),
                Inches(content.w - 0.15), Inches(row_h)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = rc(SOFT_BG)
            row_bg.line.fill.background()

            # Left accent dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(content.x + 0.25), Inches(by + row_h / 2 - 0.06),
                Inches(0.12), Inches(0.12)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rc(CARD_ACCENTS[bi % len(CARD_ACCENTS)])
            dot.line.fill.background()

            # Text
            txt = slide.shapes.add_textbox(
                Inches(content.x + 0.55), Inches(by + 0.06),
                Inches(content.w - 0.7), Inches(row_h - 0.12)
            )
            tf = txt.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = b_title
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = rc(c['text'])
            run.font.name = resolve_font(b_title, FONT_TITLE)
            if b_body:
                run2 = p.add_run()
                run2.text = ' — ' + b_body
                run2.font.size = Pt(14)
                run2.font.bold = False
                run2.font.color.rgb = rc(c['dark2'])
                run2.font.name = resolve_font(b_body, FONT_BODY)

    add_icon_to_slide(slide, spec, 'lucide:alert-triangle', c['accent4'])
    add_slide_number(slide, 3)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 4: Cards — Solution
    # =========================================================================
    slide_idx = 3
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=28)
    add_key_message(slide, spec, sd['key_message_text'], font_size=16)
    add_accent_divider(slide, spec)

    cards = spec.cards
    if cards:
        card_accent_colors = [c['accent1'], c['accent4'], c['accent6'], c['accent3']]
        card_icon_names = ['lucide:bot', 'lucide:workflow', 'lucide:database', 'lucide:globe']
        for ci, bullet in enumerate(sd['bullets'][:4]):
            parts = bullet.split(' — ', 1)
            card_title = parts[0]
            card_body = parts[1] if len(parts) > 1 else ''
            cr = cards.card_rect(ci)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cr.x), Inches(cr.y),
                Inches(cr.w), Inches(cr.h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = rc(c['white'])
            card.line.color.rgb = rc(card_accent_colors[ci])
            card.line.width = Pt(1.5)

            # Top accent stripe on card
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(cr.x), Inches(cr.y),
                Inches(cr.w), Inches(0.05)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rc(card_accent_colors[ci])
            stripe.line.fill.background()

            # Card icon
            ci_icon = fetch_icon(card_icon_names[ci], color_hex=card_accent_colors[ci])
            if ci_icon:
                safe_add_picture(slide.shapes, ci_icon,
                    Inches(cr.x + cr.w - 0.55), Inches(cr.y + 0.08),
                    width=Inches(0.4), height=Inches(0.4))

            write_panel_text(card, card_title, card_body, title_pt=14, body_pt=12)

    add_icon_to_slide(slide, spec, 'lucide:layers', c['accent1'])
    add_slide_number(slide, 4)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 5: Comparison
    # =========================================================================
    slide_idx = 4
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=26)
    add_key_message(slide, spec, sd['key_message_text'], font_size=15)
    add_accent_divider(slide, spec)

    comp = spec.comparison
    if comp:
        # Left panel — Agent Builder
        left_data = sd['bullets'][0] if len(sd['bullets']) > 0 else ''
        lp = left_data.split(' — ', 1)
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(comp.left.x), Inches(comp.left.y),
            Inches(comp.left.w), Inches(comp.left.h)
        )
        left_panel.fill.solid()
        left_panel.fill.fore_color.rgb = rc(SOFT_BG)
        left_panel.line.color.rgb = rc(c['accent5'])
        left_panel.line.width = Pt(1.5)

        # Left header bar
        lh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(comp.left.x), Inches(comp.left.y),
            Inches(comp.left.w), Inches(0.45)
        )
        lh.fill.solid()
        lh.fill.fore_color.rgb = rc(c['accent5'])
        lh.line.fill.background()
        lh_tf = lh.text_frame
        lh_tf.word_wrap = True
        lh_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        lh_tf.margin_left = Inches(0.12)
        lh_p = lh_tf.paragraphs[0]
        lh_p.alignment = PP_ALIGN.LEFT
        lh_run = lh_p.add_run()
        lh_run.text = 'M365 Agent Builder'
        lh_run.font.size = Pt(16)
        lh_run.font.bold = True
        lh_fg = ensure_contrast(c['text'], c['accent5'])
        lh_run.font.color.rgb = rc(lh_fg)
        lh_run.font.name = FONT_TITLE

        # Left body content — bullets 0 and 2
        left_body_texts = []
        if len(sd['bullets']) > 0:
            p0 = sd['bullets'][0].split(' — ', 1)
            left_body_texts.append(p0[1] if len(p0) > 1 else p0[0])
        if len(sd['bullets']) > 2:
            p2 = sd['bullets'][2].split(' — ', 1)
            left_body_texts.append(p2[1] if len(p2) > 1 else p2[0])

        lb_y = comp.left.y + 0.55
        lb_h = comp.left.h - 0.65
        left_body = slide.shapes.add_textbox(
            Inches(comp.left.x + 0.15), Inches(lb_y),
            Inches(comp.left.w - 0.3), Inches(lb_h)
        )
        lb_tf = left_body.text_frame
        lb_tf.word_wrap = True
        lb_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for ti, txt in enumerate(left_body_texts):
            p = lb_tf.paragraphs[0] if ti == 0 else lb_tf.add_paragraph()
            p.space_before = Pt(6) if ti > 0 else Pt(0)
            run = p.add_run()
            run.text = '→ ' + txt
            run.font.size = Pt(14)
            run.font.color.rgb = rc(c['dark2'])
            run.font.name = resolve_font(txt, FONT_BODY)

        # Right panel — Copilot Studio
        right_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(comp.right.x), Inches(comp.right.y),
            Inches(comp.right.w), Inches(comp.right.h)
        )
        right_panel.fill.solid()
        right_panel.fill.fore_color.rgb = rc(c['white'])
        right_panel.line.color.rgb = rc(c['accent1'])
        right_panel.line.width = Pt(2)

        # Right header bar
        rh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(comp.right.x), Inches(comp.right.y),
            Inches(comp.right.w), Inches(0.45)
        )
        rh.fill.solid()
        rh.fill.fore_color.rgb = rc(c['accent1'])
        rh.line.fill.background()
        rh_tf = rh.text_frame
        rh_tf.word_wrap = True
        rh_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        rh_tf.margin_left = Inches(0.12)
        rh_p = rh_tf.paragraphs[0]
        rh_p.alignment = PP_ALIGN.LEFT
        rh_run = rh_p.add_run()
        rh_run.text = 'Copilot Studio'
        rh_run.font.size = Pt(16)
        rh_run.font.bold = True
        rh_fg = ensure_contrast(c['text'], c['accent1'])
        rh_run.font.color.rgb = rc(rh_fg)
        rh_run.font.name = FONT_TITLE

        # Right body content — bullets 1 and 3
        right_body_texts = []
        if len(sd['bullets']) > 1:
            p1 = sd['bullets'][1].split(' — ', 1)
            right_body_texts.append(p1[1] if len(p1) > 1 else p1[0])
        if len(sd['bullets']) > 3:
            p3 = sd['bullets'][3].split(' — ', 1)
            right_body_texts.append(p3[1] if len(p3) > 1 else p3[0])

        rb_y = comp.right.y + 0.55
        rb_h = comp.right.h - 0.65
        right_body = slide.shapes.add_textbox(
            Inches(comp.right.x + 0.15), Inches(rb_y),
            Inches(comp.right.w - 0.3), Inches(rb_h)
        )
        rb_tf = right_body.text_frame
        rb_tf.word_wrap = True
        rb_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for ti, txt in enumerate(right_body_texts):
            p = rb_tf.paragraphs[0] if ti == 0 else rb_tf.add_paragraph()
            p.space_before = Pt(6) if ti > 0 else Pt(0)
            run = p.add_run()
            run.text = '→ ' + txt
            run.font.size = Pt(14)
            run.font.color.rgb = rc(c['dark2'])
            run.font.name = resolve_font(txt, FONT_BODY)

    add_icon_to_slide(slide, spec, 'lucide:git-branch', c['accent4'])
    add_slide_number(slide, 5)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 6: Timeline
    # =========================================================================
    slide_idx = 5
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=26)
    add_key_message(slide, spec, sd['key_message_text'], font_size=15)
    add_accent_divider(slide, spec)

    tl = spec.timeline
    if tl:
        # Vertical timeline line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(tl.line_x), Inches(tl.line_y),
            Inches(0.04), Inches(tl.line_h)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = rc(c['accent1'])
        line.line.fill.background()

        timeline_colors = [c['accent1'], c['accent3'], c['accent4'], c['accent5'], c['accent6'], c['accent2']]

        for ti, bullet in enumerate(sd['bullets'][:6]):
            parts = bullet.split(' — ', 1)
            date_str = parts[0]
            desc_str = parts[1] if len(parts) > 1 else ''

            nr = tl.node_rect(ti)

            # Dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(tl.dot_x), Inches(nr.y + 0.06),
                Inches(tl.dot_size), Inches(tl.dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rc(timeline_colors[ti % len(timeline_colors)])
            dot.line.fill.background()

            # Date label
            date_box = slide.shapes.add_textbox(
                Inches(nr.x), Inches(nr.y),
                Inches(1.2), Inches(nr.h)
            )
            d_tf = date_box.text_frame
            d_tf.word_wrap = False
            d_tf.auto_size = MSO_AUTO_SIZE.NONE
            d_p = d_tf.paragraphs[0]
            d_run = d_p.add_run()
            d_run.text = date_str
            d_run.font.size = Pt(12)
            d_run.font.bold = True
            d_run.font.color.rgb = rc(c['text'])
            d_run.font.name = FONT_TITLE

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(nr.x + 1.25), Inches(nr.y),
                Inches(nr.w - 1.25), Inches(nr.h)
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            desc_p = desc_tf.paragraphs[0]
            desc_run = desc_p.add_run()
            desc_run.text = desc_str
            desc_run.font.size = Pt(13)
            desc_run.font.color.rgb = rc(c['dark2'])
            desc_run.font.name = resolve_font(desc_str, FONT_BODY)

    add_icon_to_slide(slide, spec, 'lucide:calendar', c['accent6'])
    add_slide_number(slide, 6)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 7: Stats — Adoption Cost
    # =========================================================================
    slide_idx = 6
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=28)
    add_key_message(slide, spec, sd['key_message_text'], font_size=15)
    add_accent_divider(slide, spec)

    stats = spec.stats
    if stats:
        stat_labels = ['Free Trial', 'M365 Included', 'Standalone']
        stat_icons_names = ['lucide:gift', 'lucide:package', 'lucide:crown']
        stat_colors = [c['accent1'], c['accent4'], c['accent6']]
        stat_descriptions = []
        for bullet in sd['bullets'][:3]:
            parts = bullet.split(' — ', 1)
            stat_descriptions.append(parts[1] if len(parts) > 1 else parts[0])

        for si in range(3):
            br = stats.box_rect(si)
            # Card background
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(br.x), Inches(br.y),
                Inches(br.w), Inches(br.h)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = rc(c['white'])
            box.line.color.rgb = rc(stat_colors[si])
            box.line.width = Pt(2)

            # Top accent bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(br.x), Inches(br.y),
                Inches(br.w), Inches(0.06)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = rc(stat_colors[si])
            top_bar.line.fill.background()

            # Stat icon
            s_icon = fetch_icon(stat_icons_names[si], color_hex=stat_colors[si])
            if s_icon:
                safe_add_picture(slide.shapes, s_icon,
                    Inches(br.x + br.w / 2 - 0.35), Inches(br.y + 0.2),
                    width=Inches(0.7), height=Inches(0.7))

            # Stat label
            label_box = slide.shapes.add_textbox(
                Inches(br.x + 0.1), Inches(br.y + 1.0),
                Inches(br.w - 0.2), Inches(0.4)
            )
            l_tf = label_box.text_frame
            l_tf.word_wrap = True
            l_tf.auto_size = MSO_AUTO_SIZE.NONE
            l_p = l_tf.paragraphs[0]
            l_p.alignment = PP_ALIGN.CENTER
            l_run = l_p.add_run()
            l_run.text = stat_labels[si]
            l_run.font.size = Pt(18)
            l_run.font.bold = True
            l_run.font.color.rgb = rc(c['text'])
            l_run.font.name = FONT_TITLE

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(br.x + 0.12), Inches(br.y + 1.45),
                Inches(br.w - 0.24), Inches(br.h - 1.65)
            )
            d_tf = desc_box.text_frame
            d_tf.word_wrap = True
            d_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            d_p = d_tf.paragraphs[0]
            d_p.alignment = PP_ALIGN.CENTER
            d_run = d_p.add_run()
            d_run.text = stat_descriptions[si]
            d_run.font.size = Pt(13)
            d_run.font.color.rgb = rc(c['dark2'])
            d_run.font.name = resolve_font(stat_descriptions[si], FONT_BODY)

    # Footer bar with Pay-As-You-Go info
    footer = spec.footer_rect
    if footer and len(sd['bullets']) > 3:
        paygo = sd['bullets'][3]
        parts = paygo.split(' — ', 1)
        footer_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(footer.x), Inches(footer.y),
            Inches(footer.w), Inches(footer.h)
        )
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = rc(SOFT_BG)
        footer_shape.line.fill.background()
        footer_shape.name = 'footer_paygo'
        write_panel_text(footer_shape, parts[0] if len(parts) > 1 else paygo,
                         parts[1] if len(parts) > 1 else '', title_pt=14, body_pt=12)

    add_icon_to_slide(slide, spec, 'lucide:credit-card', c['accent3'])
    add_slide_number(slide, 7)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 8: Cards — Governance
    # =========================================================================
    slide_idx = 7
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=26)
    add_key_message(slide, spec, sd['key_message_text'], font_size=15)
    add_accent_divider(slide, spec)

    cards = spec.cards
    if cards:
        gov_colors = [c['accent1'], c['accent4'], c['accent5'], c['accent6']]
        gov_icons = ['lucide:shield', 'lucide:eye', 'lucide:lock', 'lucide:tag']
        for ci, bullet in enumerate(sd['bullets'][:4]):
            parts = bullet.split(' — ', 1)
            card_title = parts[0]
            card_body = parts[1] if len(parts) > 1 else ''
            cr = cards.card_rect(ci)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cr.x), Inches(cr.y),
                Inches(cr.w), Inches(cr.h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = rc(c['white'])
            card.line.color.rgb = rc(gov_colors[ci])
            card.line.width = Pt(1.5)

            # Left accent stripe
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(cr.x), Inches(cr.y),
                Inches(0.06), Inches(cr.h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rc(gov_colors[ci])
            stripe.line.fill.background()

            # Card icon
            gi = fetch_icon(gov_icons[ci], color_hex=gov_colors[ci])
            if gi:
                safe_add_picture(slide.shapes, gi,
                    Inches(cr.x + cr.w - 0.5), Inches(cr.y + cr.h / 2 - 0.2),
                    width=Inches(0.4), height=Inches(0.4))

            write_panel_text(card, card_title, card_body, title_pt=14, body_pt=12)

    add_icon_to_slide(slide, spec, 'lucide:shield-check', c['accent1'])
    add_slide_number(slide, 8)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 9: Bullets — Execution Plan
    # =========================================================================
    slide_idx = 8
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=28)
    add_key_message(slide, spec, sd['key_message_text'], font_size=16)
    add_accent_divider(slide, spec)

    content = spec.content_rect
    if content and sd.get('bullets'):
        phase_count = len(sd['bullets'])
        gap = 0.18
        row_h = (content.h - gap * (phase_count - 1)) / phase_count
        row_h = min(row_h, 1.15)
        phase_colors = [c['accent1'], c['accent4'], c['accent6']]
        phase_nums = ['01', '02', '03']

        for pi, bullet in enumerate(sd['bullets'][:3]):
            parts = bullet.split(' — ', 1)
            phase_title = parts[0]
            phase_body = parts[1] if len(parts) > 1 else ''
            py = content.y + pi * (row_h + gap)

            # Phase card
            phase_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(content.x + 0.15), Inches(py),
                Inches(content.w - 0.15), Inches(row_h)
            )
            phase_card.fill.solid()
            phase_card.fill.fore_color.rgb = rc(c['white'])
            phase_card.line.color.rgb = rc(phase_colors[pi])
            phase_card.line.width = Pt(1.5)

            # Phase number badge
            badge_w = 0.55
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(content.x + 0.25), Inches(py + row_h / 2 - 0.22),
                Inches(badge_w), Inches(0.44)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = rc(phase_colors[pi])
            badge.line.fill.background()
            b_tf = badge.text_frame
            b_tf.word_wrap = False
            b_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            b_tf.margin_left = Inches(0.02)
            b_tf.margin_right = Inches(0.02)
            b_p = b_tf.paragraphs[0]
            b_p.alignment = PP_ALIGN.CENTER
            b_run = b_p.add_run()
            b_run.text = phase_nums[pi]
            b_run.font.size = Pt(18)
            b_run.font.bold = True
            badge_fg = ensure_contrast(c['white'], phase_colors[pi])
            b_run.font.color.rgb = rc(badge_fg)
            b_run.font.name = FONT_TITLE

            # Phase text
            txt = slide.shapes.add_textbox(
                Inches(content.x + 0.95), Inches(py + 0.08),
                Inches(content.w - 1.15), Inches(row_h - 0.16)
            )
            t_tf = txt.text_frame
            t_tf.word_wrap = True
            t_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = t_tf.paragraphs[0]
            run = p.add_run()
            run.text = phase_title
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = rc(c['text'])
            run.font.name = resolve_font(phase_title, FONT_TITLE)
            if phase_body:
                p2 = t_tf.add_paragraph()
                p2.space_before = Pt(3)
                run2 = p2.add_run()
                run2.text = phase_body
                run2.font.size = Pt(13)
                run2.font.color.rgb = rc(c['dark2'])
                run2.font.name = resolve_font(phase_body, FONT_BODY)

    add_icon_to_slide(slide, spec, 'lucide:rocket', c['accent4'])
    add_slide_number(slide, 9)
    add_notes(slide, sd.get('notes', ''))

    # =========================================================================
    # SLIDE 10: Summary
    # =========================================================================
    slide_idx = 9
    spec = PRECOMPUTED_LAYOUT_SPECS[slide_idx]
    sd = layout_input[slide_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_vertical_accent_bar(slide, spec)

    add_title_textbox(slide, spec, sd['title_text'], font_size=30)
    add_key_message(slide, spec, sd['key_message_text'], font_size=16)
    add_accent_divider(slide, spec)

    # Summary box
    sb = spec.summary_box
    if sb:
        summary_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sb.x), Inches(sb.y),
            Inches(sb.w), Inches(sb.h)
        )
        summary_shape.fill.solid()
        summary_shape.fill.fore_color.rgb = rc(SOFT_BG)
        summary_shape.line.color.rgb = rc(c['accent1'])
        summary_shape.line.width = Pt(2)
        s_tf = summary_shape.text_frame
        s_tf.word_wrap = True
        s_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        s_tf.margin_left = Inches(0.15)
        s_tf.margin_right = Inches(0.12)
        s_tf.margin_top = Inches(0.1)
        s_tf.margin_bottom = Inches(0.08)
        s_p = s_tf.paragraphs[0]
        s_p.alignment = PP_ALIGN.LEFT
        s_run = s_p.add_run()
        s_run.text = sd['key_message_text']
        s_run.font.size = Pt(16)
        s_run.font.bold = True
        s_run.font.color.rgb = rc(c['text'])
        s_run.font.name = resolve_font(sd['key_message_text'], FONT_TITLE)

    # Action items
    content = spec.content_rect
    if content and sd.get('bullets'):
        action_count = len(sd['bullets'])
        gap = 0.15
        row_h = min((content.h - gap * (action_count - 1)) / action_count, 0.85)
        action_colors = [c['accent1'], c['accent4'], c['accent6']]

        for ai, bullet in enumerate(sd['bullets'][:3]):
            ay = content.y + ai * (row_h + gap)

            # Action row
            action_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(content.x + 0.15), Inches(ay),
                Inches(content.w - 0.15), Inches(row_h)
            )
            action_bg.fill.solid()
            action_bg.fill.fore_color.rgb = rc(c['white'])
            action_bg.line.color.rgb = rc(action_colors[ai])
            action_bg.line.width = Pt(1.5)

            # Left accent
            la = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(content.x + 0.15), Inches(ay),
                Inches(0.06), Inches(row_h)
            )
            la.fill.solid()
            la.fill.fore_color.rgb = rc(action_colors[ai])
            la.line.fill.background()

            # Text
            txt = slide.shapes.add_textbox(
                Inches(content.x + 0.4), Inches(ay + 0.08),
                Inches(content.w - 0.6), Inches(row_h - 0.16)
            )
            t_tf = txt.text_frame
            t_tf.word_wrap = True
            t_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = t_tf.paragraphs[0]
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(16)
            run.font.color.rgb = rc(c['text'])
            run.font.name = resolve_font(bullet, FONT_BODY)

    # Image in icon area
    img_path_10 = os.path.join(IMAGES_DIR,
        '10-start-today-three-immediate-actions-839cb11b5f.jpg')
    ir = spec.icon_rect
    if ir:
        safe_add_picture(slide.shapes, img_path_10,
            Inches(ir.x), Inches(ir.y),
            width=Inches(ir.w), height=Inches(ir.h))

    add_icon_to_slide(slide, spec, 'lucide:arrow-right', c['accent4'])
    add_slide_number(slide, 10)
    add_notes(slide, sd.get('notes', ''))

    # Save
    prs.save(output_path)